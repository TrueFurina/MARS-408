# -*- coding: utf-8 -*-
"""火山杯演示 PPT 生成脚本 — 数字口径与 README_火山杯-最终版.md 一致（2026-08-28 定稿）"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── 主题色 ──
NAVY   = RGBColor(0x12, 0x2A, 0x4A)   # 深蓝主色
BLUE   = RGBColor(0x1F, 0x5C, 0x99)   # 中蓝
CYAN   = RGBColor(0x00, 0x9E, 0xA8)   # 青色点缀
LIGHT  = RGBColor(0xF2, 0xF6, 0xFB)   # 浅底
GRAY   = RGBColor(0x5A, 0x6B, 0x7B)   # 次级文字
DARK   = RGBColor(0x1A, 0x24, 0x2F)   # 正文
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xE8, 0x6A, 0x17)

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]

FONT = "Microsoft YaHei"

def _set_ea(r, font=FONT):
    """设置中文字体（a:ea / a:cs）"""
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", font)

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def text(slide, x, y, w, h, runs, size=16, color=DARK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.0, font=FONT):
    """runs: str 或 [(text, dict_override), ...]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [(runs, {})]
    for i, (t, ov) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        p.line_spacing = ov.get("line_spacing", line_spacing)
        r = p.add_run()
        r.text = t
        r.font.size = Pt(ov.get("size", size))
        r.font.bold = ov.get("bold", bold)
        r.font.color.rgb = ov.get("color", color)
        r.font.name = ov.get("font", font)
        _set_ea(r)
    return tb

def header(slide, title, subtitle=None, page=None):
    rect(slide, 0, 0, SW, Inches(1.0), NAVY)
    rect(slide, 0, Inches(1.0), SW, Pt(3), CYAN)
    text(slide, Inches(0.55), Inches(0.16), Inches(9.5), Inches(0.7),
         [(title, {"size": 26, "bold": True, "color": WHITE})], anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        text(slide, Inches(9.8), Inches(0.16), Inches(3.0), Inches(0.7),
             [(subtitle, {"size": 12, "color": RGBColor(0xBF,0xD4,0xE8)})],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    if page:
        text(slide, Inches(12.45), Inches(7.05), Inches(0.7), Inches(0.35),
             [(str(page), {"size": 10, "color": GRAY})], align=PP_ALIGN.RIGHT)
    text(slide, Inches(0.45), Inches(7.05), Inches(6), Inches(0.35),
         [("MARS-408 · 2026 福建高校「火山杯」Agent 创新大赛", {"size": 9, "color": GRAY})])

def bullet_block(slide, x, y, w, h, items, size=15, gap=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r1 = p.add_run(); r1.text = "▍"
        r1.font.size = Pt(size); r1.font.color.rgb = CYAN; r1.font.name = FONT
        _set_ea(r1)
        r2 = p.add_run(); r2.text = it
        r2.font.size = Pt(size); r2.font.color.rgb = DARK; r2.font.name = FONT
        _set_ea(r2)
    return tb

def card(slide, x, y, w, h, title, body, tcolor=BLUE, tsize=16, bsize=12.5):
    rect(slide, x, y, w, h, LIGHT)
    rect(slide, x, y, Inches(0.06), h, tcolor)
    text(slide, x + Inches(0.22), y + Inches(0.14), w - Inches(0.4), Inches(0.4),
         [(title, {"size": tsize, "bold": True, "color": tcolor})])
    text(slide, x + Inches(0.22), y + Inches(0.55), w - Inches(0.4), h - Inches(0.7),
         [(body, {"size": bsize, "color": DARK})], line_spacing=1.15)

def flow_boxes(slide, y, items, x0=Inches(0.55), w_total=Inches(12.25)):
    """items: [(title, sub, color)] 横向流程块"""
    n = len(items)
    gap = Inches(0.25)
    bw = Emu(int((w_total - gap * (n - 1)) / n))
    x = x0
    for i, (t, s, c) in enumerate(items):
        rect(slide, x, y, bw, Inches(1.55), c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(slide, x + Inches(0.15), y + Inches(0.2), bw - Inches(0.3), Inches(0.55),
             [(t, {"size": 15, "bold": True, "color": WHITE})], align=PP_ALIGN.CENTER)
        text(slide, x + Inches(0.15), y + Inches(0.75), bw - Inches(0.3), Inches(0.7),
             [(s, {"size": 11, "color": RGBColor(0xE8,0xF1,0xFA)})], align=PP_ALIGN.CENTER, line_spacing=1.1)
        if i < n - 1:
            ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + bw + Inches(0.02), y + Inches(0.55), Inches(0.21), Inches(0.45))
            ar.fill.solid(); ar.fill.fore_color.rgb = CYAN; ar.line.fill.background(); ar.shadow.inherit = False
        x = Emu(int(x + bw + gap))

def table(slide, x, y, w, rows, col_widths, header_fill=NAVY, size=13, row_h=0.42):
    n_r, n_c = len(rows), len(rows[0])
    h_in = row_h * n_r
    gt = slide.shapes.add_table(n_r, n_c, x, y, w, Inches(h_in)).table
    for j, cw in enumerate(col_widths):
        gt.columns[j].width = Inches(cw)
    for i in range(n_r):
        gt.rows[i].height = Inches(row_h)
        for j in range(n_c):
            cell = gt.cell(i, j)
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = rows[i][j]
            r.font.name = FONT; _set_ea(r)
            r.font.size = Pt(size)
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                r.font.color.rgb = WHITE; r.font.bold = True
                p.alignment = PP_ALIGN.CENTER
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
                r.font.color.rgb = DARK
                if j == 0:
                    r.font.bold = True
    return gt

# ══════════════ 1 封面 ══════════════
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.9), SW, Pt(3), CYAN)
text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.5),
     [("2026 福建高校「火山杯」Agent 创新大赛", {"size": 18, "color": CYAN, "bold": True})])
text(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.5),
     [("MARS-408", {"size": 54, "bold": True, "color": WHITE})])
text(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.8),
     [("考研多智能体个性化学习系统", {"size": 30, "bold": True, "color": WHITE})])
text(s, Inches(0.9), Inches(4.35), Inches(11.5), Inches(0.5),
     [("让 AI 从「回答问题」到「真正懂你」—— 10 节点多智能体流水线驱动的考研学习教练", {"size": 15, "color": RGBColor(0xC9,0xDA,0xEC)})])
text(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(0.9),
     [("覆盖 408 计算机考研四科：数据结构 · 计算机组成原理 · 操作系统 · 计算机网络\n赛道方向：未来学习中心（个性化学习规划 / 知识图谱 / 学习效果追踪）", {"size": 13, "color": RGBColor(0xAF,0xC3,0xDA), "line_spacing": 1.3})])

# ══════════════ 2 痛点与场景 ══════════════
s = add_slide()
header(s, "一、痛点：为什么传统 AI 助教不够用", "未来学习中心 · 场景切入", 2)
text(s, Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.5),
     [("408 计算机考研是百万级考研大军中竞争最激烈的赛道，备考长期面临三大困境：", {"size": 15, "color": DARK})])
cards = [
    ("困境 1 · 资料海量", "教材、题库、网课、笔记信息过载，\n「该学什么」全靠个人感觉", BLUE),
    ("困境 2 · 无人诊断", "薄弱点靠猜：学没学会、错在哪、\n下一步学什么，没有客观依据", ORANGE),
    ("困境 3 · 学练割裂", "刷题与讲解脱节，错题无人讲解，\n讲解后无人出题验证", CYAN),
]
x = Inches(0.55); w = Inches(3.95)
for i, (t, b, c) in enumerate(cards):
    card(s, Emu(int(x + i * (w + Inches(0.2)))), Inches(1.95), w, Inches(2.15), t, b, tcolor=c)
rect(s, Inches(0.55), Inches(4.35), Inches(12.25), Inches(2.1), LIGHT)
rect(s, Inches(0.55), Inches(4.35), Inches(0.06), Inches(2.1), CYAN)
text(s, Inches(0.85), Inches(4.55), Inches(11.7), Inches(1.8),
     [("传统 AI 助手只能「一问一答」：", {"size": 15, "bold": True, "color": NAVY}),
      ("回答完就结束，没有诊断、没有规划、没有追踪，无法完成「诊断 → 规划 → 讲解 → 练习 → 复盘」的学习闭环。", {"size": 14, "color": DARK}),
      ("\nMARS-408 的目标：让 AI 从「回答问题的工具」升级为「真正懂你的学习教练」。", {"size": 15, "bold": True, "color": ORANGE})], line_spacing=1.25)

# ══════════════ 3 产品定位 ══════════════
s = add_slide()
header(s, "二、产品定位：多智能体学习教练", "", 3)
flow_boxes(s, Inches(1.75), [
    ("学情诊断", "diagnostician\n定位薄弱知识点", BLUE),
    ("路径规划", "planner / path_planner\n制定个性化学习路径", NAVY),
    ("讲解生成", "generator_cluster\n7 种资源并行生成", CYAN),
    ("出题评估", "assessor\n按难度出题并批改", ORANGE),
    ("复盘迭代", "critic + evidence_check\n质量校验与证据核查", BLUE),
])
text(s, Inches(0.55), Inches(3.7), Inches(12.2), Inches(0.5),
     [("完整闭环：每一次提问都不是终点，而是下一次更懂你的起点。", {"size": 14, "bold": True, "color": NAVY})])
cards = [
    ("主动规划", "系统自动拆解学习任务、\n规划分步路径，而非被动应答", BLUE),
    ("多轮交互", "支持追问、引导与纠错，\n对话上下文持续积累", CYAN),
    ("工具协同", "知识库检索 + 计算工具 + 多模态生成，\nAgent 不只是「说」而是「做」", ORANGE),
]
x = Inches(0.55); w = Inches(3.95)
for i, (t, b, c) in enumerate(cards):
    card(s, Emu(int(x + i * (w + Inches(0.2)))), Inches(4.35), w, Inches(2.0), t, b, tcolor=c)

# ══════════════ 4 核心创新 · 防幻觉 ══════════════
s = add_slide()
header(s, "三、核心创新 ①：多智能体共识，从机制上防幻觉", "区别于普通问答 / 简单内容生成", 4)
text(s, Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.5),
     [("为什么是「多个 Agent」而不是「一个大模型」？—— 单模型自说自话无法自我纠错：", {"size": 15, "color": DARK})])
flow_boxes(s, Inches(2.0), [
    ("多 Agent 独立作答", "不同角色各出答案\n（讲解 / 出题 / 校验）", BLUE),
    ("GoMARL 加权共识", "基于 GoMARL 思想的\n加权共识引擎裁决分歧", NAVY),
    ("冲突消解", "发现知识矛盾（如\n「三次握手 vs 四次挥手」）", ORANGE),
    ("证据链复核", "知识库检索证据 +\n真实大模型复核事实", CYAN),
])
card(s, Inches(0.55), Inches(3.95), Inches(12.25), Inches(2.4), "冲突消解引擎 · 工作流程",
     "① 多 Agent 输出不一致 → ② 冲突消解引擎基于证据链检索知识库 → ③ 由真实大模型复核事实后裁决 → ④ 输出带证据标注、可信度评分的最终答案。\n\n与「提示词约束防幻觉」的本质区别：分歧被显式暴露、证据被显式检索、事实被显式复核——AI 幻觉从机制上被拦截。",
     tcolor=NAVY, tsize=16, bsize=13)

# ══════════════ 5 核心创新 · 个性化闭环 ══════════════
s = add_slide()
header(s, "三、核心创新 ②：三层画像 + 个性化闭环", "", 5)
cards = [
    ("短期 · 对话记忆", "当前会话上下文持续跟踪，\n追问与纠错不丢失", BLUE),
    ("中期 · 语义记忆", "知识点掌握度动态更新，\n薄弱点实时感知", CYAN),
    ("长期 · 成长画像", "学习轨迹持续积累，\n越用越懂你", ORANGE),
]
x = Inches(0.55); w = Inches(3.95)
for i, (t, b, c) in enumerate(cards):
    card(s, Emu(int(x + i * (w + Inches(0.2)))), Inches(1.6), w, Inches(1.9), t, b, tcolor=c)
text(s, Inches(0.55), Inches(3.75), Inches(12.2), Inches(0.5),
     [("画像驱动 · 精准到每个知识点：", {"size": 15, "bold": True, "color": NAVY})])
bullet_block(s, Inches(0.85), Inches(4.3), Inches(11.5), Inches(2.3), [
    "个性化学习路径：依据画像与薄弱点动态规划「下一步学什么」，薄弱点优先",
    "专属练习生成：按科目 / 章节 / 难度自动出题，答后自动批改并反哺画像",
    "学习效果追踪：多维评估报告（知识点掌握 / 答题正确率 / 薄弱点变化），教师端支持班级学情聚合",
], size=14)

# ══════════════ 6 技术架构 ══════════════
s = add_slide()
header(s, "四、技术实现 · 整体架构", "", 6)
layers = [
    ("前端层 · Vue 3 + TypeScript", "68 个页面 · Vite 构建 · 多端多角色（学生端 / 教师看板）", BLUE),
    ("智能体层 · FastAPI + LangGraph", "10 节点多智能体流水线 · GoMARL 共识引擎 · FrugalRAG 检索引擎", NAVY),
    ("模型层 · 双通道大模型自动容灾", "讯飞星火 X2（主通道） + DeepSeek（兜底） · 调用失败自动切换", CYAN),
    ("数据层 · 多级存储与自动降级", "Milvus / InMemoryVectorStore · PostgreSQL / SQLite · Redis / 内存", ORANGE),
]
y = Inches(1.45); lh = Inches(1.15); gap = Inches(0.22)
for i, (t, b, c) in enumerate(layers):
    yy = Emu(int(y + i * (lh + gap)))
    rect(s, Inches(0.55), yy, Inches(12.25), lh, LIGHT)
    rect(s, Inches(0.55), yy, Inches(0.06), lh, c)
    text(s, Inches(0.85), Emu(int(yy + Inches(0.12))), Inches(11.7), Inches(0.4),
         [(t, {"size": 16, "bold": True, "color": c})])
    text(s, Inches(0.85), Emu(int(yy + Inches(0.55))), Inches(11.7), Inches(0.45),
         [(b, {"size": 13, "color": DARK})])
    if i < len(layers) - 1:
        ar = slide_arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.5), Emu(int(yy + lh + Inches(0.01))), Inches(0.35), Inches(0.2))
        ar.fill.solid(); ar.fill.fore_color.rgb = CYAN; ar.line.fill.background(); ar.shadow.inherit = False

# ══════════════ 7 10 节点流水线 ══════════════
s = add_slide()
header(s, "四、技术实现 · 10 节点多智能体流水线", "LangGraph 状态图编排", 7)
rows = [
    ["节点", "职责", "产出"],
    ["coordinator", "识别意图、全局协调、任务分派", "任务路由"],
    ["diagnostician", "学情诊断，定位薄弱知识点", "诊断报告"],
    ["planner / path_planner", "分析画像、制定分步学习计划", "学习路径"],
    ["retriever", "向量 + BM25 混合检索知识库", "证据片段"],
    ["generator_cluster", "7 种资源并行生成", "讲解 / 习题 / 导图 / PPT / 视频脚本等"],
    ["assessor", "按知识点与难度出题、批改反馈", "练习题 / 评分"],
    ["critic / evidence_check / quality_gate", "质量审核 · 证据核查 · 产物验收", "审核报告 / 验收结论"],
]
table(s, Inches(0.55), Inches(1.35), Inches(12.25), rows, [3.4, 5.55, 3.3], size=12.5, row_h=0.62)
text(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.5),
     [("支持 7 种学习资源并行生成：讲解文档 · 练习题 · 思维导图 · 拓展阅读 · PPT 大纲 · 代码实操 · 视频脚本（SSE 流式进度推送）", {"size": 12.5, "color": GRAY})])

# ══════════════ 8 知识库与检索 ══════════════
s = add_slide()
header(s, "四、技术实现 · 知识库与 FrugalRAG 检索", "", 8)
rows = [
    ["数据", "数量", "说明"],
    ["知识分片", "1883", "739 knowledge_point + 1144 knowledge_variant，覆盖四科高频考点"],
    ["练习题", "200", "选择 / 填空 / 简答，覆盖四科"],
    ["向量条目", "2083", "全部真实 E5 嵌入（768 维），0 零向量"],
    ["知识群组", "26", "按四科章节划分，供跨群冲突检测"],
]
table(s, Inches(0.55), Inches(1.35), Inches(12.25), rows, [2.2, 1.6, 8.45], size=12.5, row_h=0.55)
text(s, Inches(0.55), Inches(4.0), Inches(12.2), Inches(0.5),
     [("检索链路：", {"size": 14, "bold": True, "color": NAVY}),
      ("提问 → E5 向量检索 → BM25 全文检索 → 融合排序 → 个性化重排 → 增强生成", {"size": 14, "color": DARK})])
card(s, Inches(0.55), Inches(4.6), Inches(12.25), Inches(1.9), "工程韧性：检索降级守卫",
     "检索链路异常时自动降级 BM25-only（携带 _degraded 标记），绝不静默返回空结果——保证演示与真实使用场景下系统始终可用、结果始终可解释。",
     tcolor=CYAN, tsize=14.5, bsize=12.5)

# ══════════════ 9 量化实测 ══════════════
s = add_slide()
header(s, "五、量化实测 · 全部可复现", "真实运行产出 · 非虚构用户实验", 9)
rows = [
    ["指标", "结果", "说明"],
    ["检索增强效果", "Recall@5 +15.8% / MRR +16.0%", "对照无重排基线，CPU 真实运行（2026-08 实测）"],
    ["检索层可回答率基线", "answerable_rate 0.533", "30 题四科验证集（eval_gold），持续回归用"],
    ["自动化测试", "500+ 项", "API 契约 / 闭环流程 / 降级路径全覆盖"],
    ["LLM 容灾", "双通道自动切换", "讯飞星火 X2 主通道 → DeepSeek 兜底，超时/429 自动降级"],
]
table(s, Inches(0.55), Inches(1.5), Inches(12.25), rows, [2.9, 3.3, 6.05], size=13, row_h=0.6)
text(s, Inches(0.55), Inches(4.4), Inches(12.2), Inches(0.5),
     [("评测脚本随源码提供（py-server/experiments/），所有指标一键复现。", {"size": 14, "bold": True, "color": NAVY})])
card(s, Inches(0.55), Inches(5.0), Inches(12.25), Inches(1.55), "为什么可信",
     "所有成效数字均为真实运行产出：检索基准（benchmark_2026-08）与验证集评测（eval_gold）均有磁盘证据与复现脚本，不包含虚构的用户实验数据。",
     tcolor=ORANGE, tsize=14.5, bsize=12.5)

# ══════════════ 10 工程指标 ══════════════
s = add_slide()
header(s, "五、工程化水平 · 完整度与可展示性", "", 10)
cards = [
    ("68", "前端页面\nVue 3 + TypeScript", BLUE),
    ("170+", "API 端点\nFastAPI", CYAN),
    ("10", "Agent 节点\nLangGraph 编排", NAVY),
    ("500+", "自动化测试\n全部通过", ORANGE),
    ("双通道", "LLM 自动容灾\n星火 X2 + DeepSeek", BLUE),
    ("单机可跑", "三级组件降级\nMilvus / PG / Redis", CYAN),
]
x0 = Inches(0.55); w = Inches(1.98); gap = Inches(0.12)
for i, (num, lab, c) in enumerate(cards):
    xx = Emu(int(x0 + i * (w + gap)))
    rect(s, xx, Inches(1.5), w, Inches(2.3), LIGHT)
    rect(s, xx, Inches(1.5), w, Inches(0.06), c)
    text(s, xx, Inches(1.75), w, Inches(0.9), [(num, {"size": 30, "bold": True, "color": c})], align=PP_ALIGN.CENTER)
    text(s, xx, Inches(2.7), w, Inches(1.0), [(lab, {"size": 11.5, "color": DARK})], align=PP_ALIGN.CENTER, line_spacing=1.2)
bullet_block(s, Inches(0.85), Inches(4.2), Inches(11.5), Inches(2.5), [
    "工程完整度：API 契约测试 / 闭环流程测试 / 降级路径测试全覆盖，缺失外部组件时自动降级，保证稳定运行",
    "可展示性：SSE 流式输出 + Agent 实时进度推送，演示效果直观流畅",
    "演示视频与演示账号随提交物提供（demo / demo123456），开箱即用",
], size=14)

# ══════════════ 11 落地价值 ══════════════
s = add_slide()
header(s, "六、落地价值 · 面向真实教育场景", "", 11)
cards = [
    ("学生 · 考研刚需", "408 备考人群基数大、需求刚性：\n诊断薄弱点 → 定制路径 → 学练闭环", BLUE),
    ("教师 · 学情管理", "教师看板：班级进度 / 掌握度 / 薄弱点聚合，\n从「凭经验」到「凭数据」", CYAN),
    ("院校 · 可推广", "工程完整、单机可部署、双通道容灾，\n可直接投入高校 / 培训机构使用", ORANGE),
]
x = Inches(0.55); w = Inches(3.95)
for i, (t, b, c) in enumerate(cards):
    card(s, Emu(int(x + i * (w + Inches(0.2)))), Inches(1.6), w, Inches(2.3), t, b, tcolor=c)
text(s, Inches(0.55), Inches(4.15), Inches(12.2), Inches(0.5),
     [("赛道定位：精准命中「未来学习中心」三大关键词", {"size": 15, "bold": True, "color": NAVY})])
flow_boxes(s, Inches(4.75), [
    ("个性化学习规划", "画像驱动路径\n薄弱点优先", BLUE),
    ("知识图谱", "26 大知识群组\n四科关联视图", NAVY),
    ("学习效果追踪", "评估报告 + 教师看板\n持续反馈闭环", CYAN),
])

# ══════════════ 12 合规与致谢 ══════════════
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
text(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.6),
     [("开发工具合规声明", {"size": 28, "bold": True, "color": WHITE})])
text(s, Inches(0.9), Inches(1.95), Inches(11.5), Inches(1.6),
     [("本项目为真实可运行的代码工程（Vue 3 + TypeScript 前端 / FastAPI + LangGraph 后端），采用 Trae（字节跳动 AI IDE）作为开发工具完成核心模块的开发与迭代；代码仓库可在 Trae 中直接打开、构建并运行，满足大赛「基于 Trae 开发」的工具要求。",
       {"size": 15, "color": RGBColor(0xE8,0xF1,0xFA), "line_spacing": 1.4})])
text(s, Inches(0.9), Inches(3.55), Inches(11.5), Inches(0.5),
     [("关键可留痕模块：py-server/agents/graph.py · engines/frugal_rag.py · engines/gomarl.py · src/", {"size": 12.5, "color": RGBColor(0xAF,0xC3,0xDA)})])
rect(s, 0, Inches(4.5), SW, Pt(2), CYAN)
text(s, Inches(0.9), Inches(4.85), Inches(11.5), Inches(0.5),
     [("致谢", {"size": 20, "bold": True, "color": CYAN})])
text(s, Inches(0.9), Inches(5.45), Inches(11.5), Inches(1.2),
     [("感谢 2026 福建高校「火山杯」Agent 创新大赛主办方与火山引擎技术支持。\n开源组件：LangGraph · FastAPI · Milvus · Vue 3 · E5 等（许可清单见 OPENSOURCE_LICENSES.md）",
       {"size": 13, "color": RGBColor(0xC9,0xDA,0xEC), "line_spacing": 1.4})])
text(s, Inches(0.9), Inches(6.9), Inches(11.5), Inches(0.4),
     [("MARS-408 · 让每一次学习都有迹可循", {"size": 14, "bold": True, "color": WHITE})])

OUT = r"E:/Program/MARL/study-help-pro/deliverables/火山杯-演示PPT-最终版.pptx"
prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
