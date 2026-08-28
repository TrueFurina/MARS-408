# ============================================================
# PPT 文件构建器 (PPT File Builder)
# 将 LLM 生成的 Markdown 大纲转换为真实的 .pptx 二进制文件
# （赛题明确要求"多模态"资源，原系统只产出大纲文本/结构，无真实文件）
#
# 设计：
#   - 不依赖 LLM，纯解析 Markdown 大纲 -> 幻灯片，确定性、零额外配额
#   - 产出标准 Office Open XML (.pptx)，可用 PowerPoint / WPS 直接打开
#   - 文件落盘到 media/ppt/，通过 /media/ppt/<file>.pptx 静态服务下载
# ============================================================

import os
import re
import hashlib
import logging
from typing import Optional

logger = logging.getLogger("netlearn.ppt_builder")

MEDIA_PPT_DIR = os.path.join(os.path.dirname(__file__), "..", "media", "ppt")

# 品牌配色（与前端玻璃态设计系统一致）
_BRAND_HEX = "2D6CDF"        # 主蓝
_BRAND_SUB_HEX = "6B7280"    # 副灰


def _safe_filename(topic: str) -> str:
    # 仅保留 ASCII 字母数字与下划线，避免 URL / 下载文件名编码问题
    s = re.sub(r"[^A-Za-z0-9_]", "_", topic or "topic").strip("_")
    return (s[:40] or "mars408")


def _parse_outline_to_slides(markdown_text: str):
    """将 Markdown 大纲解析为 [(title, [bullets]), ...] 列表。

    规则：
      #  / ##  -> 新幻灯片标题（首个作为封面）
      ###      -> 作为要点标题行
      -/*/+    -> 无序要点
      1. 2.    -> 有序要点
      其它行    -> 归入当前幻灯片要点 / 或作为默认标题
    """
    slides = []
    current_title = None
    current_bullets = []

    def _flush():
        nonlocal current_title, current_bullets
        if current_title is not None or current_bullets:
            slides.append((current_title, current_bullets))
        current_title = None
        current_bullets = []

    for raw in (markdown_text or "").splitlines():
        line = raw.strip()
        if not line:
            continue
        if line.startswith("---PPT_START---") or line.startswith("---"):
            continue
        if line.startswith("# "):
            _flush()
            current_title = line[2:].strip()
        elif line.startswith("## "):
            _flush()
            current_title = line[3:].strip()
        elif line.startswith("### "):
            current_bullets.append(line[4:].strip())
        elif line.startswith(("- ", "* ", "+ ")):
            current_bullets.append(line[2:].strip())
        elif re.match(r"^\d+[.、]",
                      line) and len(line) > 2:
            current_bullets.append(re.split(r"^\d+[.、]\s*", line, maxsplit=1)[1].strip())
        else:
            if current_title is None:
                current_title = line[:40]
            else:
                current_bullets.append(line)
    _flush()
    return slides


def build_pptx(topic: str, outline_markdown: str, profile: Optional[dict] = None,
               memory_context: str = "") -> dict:
    """根据 Markdown 大纲生成真实 .pptx 文件（可选注入三层学情记忆做封面个性化）。

    Returns:
        {
            "ok": bool,
            "filename": str, "path": str, "url": str,
            "slide_count": int, "topic": str,
            "error": str (仅失败时)
        }
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        logger.error("python-pptx 未安装，无法生成 .pptx")
        return {"ok": False, "error": "python-pptx 未安装", "topic": topic}

    try:
        os.makedirs(MEDIA_PPT_DIR, exist_ok=True)

        slides = _parse_outline_to_slides(outline_markdown)
        if not slides:
            slides = [(topic or "学习主题", ["系统生成的 PPT 大纲暂不可用，请重试"])]

        prs = Presentation()

        # ── 封面（记忆薄弱点驱动个性化副标题） ──
        cover = prs.slides.add_slide(prs.slide_layouts[0])
        cover.shapes.title.text = (topic or "个性化学习资源")
        if len(cover.placeholders) > 1:
            sub = cover.placeholders[1]
            sub.text = "MARS-408 · 多智能体个性化学习系统（多模态资源生成）"
            try:
                sub.text_frame.paragraphs[0].font.size = Pt(16)
                sub.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(_BRAND_SUB_HEX)
            except Exception:
                pass
            # L1/L2/L3 三层学情记忆（低侵入：封面标注重点复习薄弱点）
            if memory_context and memory_context != "【学生记忆】暂无历史学习数据":
                import re as _re
                weak_block = _re.search(r"薄弱[：:]\s*(.+?)(?:\n|$)", memory_context)
                if weak_block:
                    weak_terms = weak_block.group(1).strip()[:40]
                    try:
                        p2 = sub.text_frame.add_paragraph()
                        p2.text = f"重点复习：{weak_terms}"
                        p2.font.size = Pt(12)
                    except Exception:
                        pass

        # ── 内容幻灯片 ──
        content_layout = prs.slide_layouts[1]
        for title, bullets in slides:
            title = (title or topic or "学习主题")[:60]
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = title
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.word_wrap = True
            # 每页最多 9 条，避免溢出
            for i, b in enumerate(bullets[:9]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = "• " + str(b)
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor.from_string("1F2937")

        # ── 落盘 ──
        digest = hashlib.sha256((outline_markdown or topic).encode("utf-8")).hexdigest()[:8]
        filename = f"{_safe_filename(topic)}_{digest}.pptx"
        path = os.path.join(MEDIA_PPT_DIR, filename)
        prs.save(path)

        result = {
            "ok": True,
            "filename": filename,
            "path": path,
            "url": f"/media/ppt/{filename}",
            "slide_count": len(prs.slides._sldIdLst),
            "topic": topic,
        }
        logger.info(f"[PPTBuilder] 生成成功: {filename} ({result['slide_count']}页)")
        return result

    except Exception as e:
        logger.exception(f"[PPTBuilder] 生成失败: {e}")
        return {"ok": False, "error": str(e), "topic": topic}
