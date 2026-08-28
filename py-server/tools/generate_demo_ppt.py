#!/usr/bin/env python3
"""MARS-408 软件杯演示 PPT 自动生成脚本

用法:
    cd py-server
    python tools/generate_demo_ppt.py
"""

import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色方案 ──
C_PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)     # 深蓝黑
C_SECONDARY = RGBColor(0x16, 0x21, 0x3E)    # 深蓝
C_ACCENT = RGBColor(0x00, 0xD2, 0xFF)       # 亮青
C_ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)      # 紫色
C_ACCENT3 = RGBColor(0x10, 0xB9, 0x81)      # 绿色
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xCC, 0xDD, 0xEE)
C_GRAY = RGBColor(0x88, 0x99, 0xAA)


def _bg(slide, color=C_PRIMARY):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, color):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text(slide, left, top, width, height, text, font_size=18, color=C_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_text(slide, left, top, width, height, items, font_size=16, color=C_LIGHT, spacing=Pt(8)):
    """添加多行要点"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = spacing
    return txBox


# ════════════════════════════════════════════════════════════════
# Slide 1: 封面
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_bg(slide)
_add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), C_ACCENT)
_add_shape(slide, Inches(0), Inches(5.8), Inches(13.333), Inches(0.04), C_ACCENT)
_add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
          "MARS-408", font_size=60, color=C_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
_add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
          "基于 GOMARL 与 FrugalRAG 的 408 考研个性化学习多智能体系统",
          font_size=24, color=C_WHITE, alignment=PP_ALIGN.CENTER)
_add_text(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.5),
          "2026 软件杯 · 作品演示", font_size=20, color=C_ACCENT, alignment=PP_ALIGN.CENTER)
_add_text(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
          "10 节点多智能体流水线  |  GOMARL 共识引擎  |  FrugalRAG 检索  |  7 种资源并行生成",
          font_size=14, color=C_GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# Slide 2: 目录
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_bg(slide)
_add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
          "目 录", font_size=36, color=C_ACCENT, bold=True)
_add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), C_ACCENT)

toc = [
    ("01", "项目背景与目标"),
    ("02", "系统架构总览"),
    ("03", "10 节点 Agent 流水线"),
    ("04", "GOMARL 共识机制"),
    ("05", "FrugalRAG 检索策略"),
    ("06", "7 种资源并行生成"),
    ("07", "NeuralMixer 神经网络"),
    ("08", "演示与效果"),
]
for i, (num, title) in enumerate(toc):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 5.8)
    y = Inches(1.8 + row * 1.3)
    _add_shape(slide, x, y, Inches(0.6), Inches(0.6), C_ACCENT2 if i % 2 == 0 else C_ACCENT3)
    _add_text(slide, x + Inches(0.05), y + Inches(0.05), Inches(0.5), Inches(0.5),
              num, font_size=20, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text(slide, x + Inches(0.8), y + Inches(0.1), Inches(4.5), Inches(0.5),
              title, font_size=20, color=C_WHITE)

# ════════════════════════════════════════════════════════════════
# Slide 3: 项目背景
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_bg(slide)
_add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
          "01  项目背景与目标", font_size=32, color=C_ACCENT, bold=True)
_add_shape(slide, Inches(0.8), Inches(1.2), Inches(12), Inches(0.03), C_ACCENT)

_add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
          "408 考研的三大痛点", font_size=22, color=C_ACCENT, bold=True)
_add_bullet_text(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(3.5), [
    "知识体系庞大：四门课程、数百个知识点",
    "个性化需求高：每位学生基础和目标不同",
    "资源质量参差不齐：难找到精准匹配的资源",
])

_add_text(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.5),
          "MARS-408 方案", font_size=22, color=C_ACCENT2, bold=True)
_add_bullet_text(slide, Inches(6.8), Inches(2.2), Inches(5.5), Inches(3.5), [
    "多智能体协同：13 个 Agent 各司其职",
    "GOMARL 共识：质量保障与冲突消解",
    "FrugalRAG：高效精准的知识检索",
    "7 种资源并行生成：个性化学习包",
])

# ════════════════════════════════════════════════════════════════
# Slide 4: 系统架构
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_bg(slide)
_add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
          "02  系统架构总览", font_size=32, color=C_ACCENT, bold=True)
_add_shape(slide, Inches(0.8), Inches(1.2), Inches(12), Inches(0.03), C_ACCENT)

# 三层架构
layers = [
    ("API 层", "89 个路由端点  |  97.8% 认证覆盖率  |  安全响应头", C_ACCENT2, Inches(0.8)),
    ("Agent 层", "10 节点 LangGraph 流水线  |  GOMARL 共识  |  FrugalRAG 检索", C_ACCENT, Inches(2.8)),
    ("数据层", "E5 向量库  |  PostgreSQL  |  Redis  |  Milvus 抽象层", C_ACCENT3, Inches(4.8)),
]
for name, desc, color, y in layers:
    _add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.5), color)
    _add_shape(slide, Inches(0.8), y, Inches(0.08), Inches(1.5), color)
    _add_text(slide, Inches(1.2), y + Inches(0.2), Inches(4), Inches(0.5),
              name, font_size=22, color=color, bold=True)
    _add_text(slide, Inches(1.2), y + Inches(0.8), Inches(10), Inches(0.5),
              desc, font_size=16, color=C_LIGHT)

# 技术栈标签
techs = ["Vue 3 + TypeScript", "FastAPI + LangGraph", "GOMARL", "FrugalRAG", "E5 + BM25", "PyTorch"]
for i, t in enumerate(techs):
    x = Inches(0.8 + i * 2)
    _add_shape(slide, x, Inches(6.8), Inches(1.8), Inches(0.5), C_ACCENT2 if i % 2 == 0 else C_ACCENT3)
    _add_text(slide, x, Inches(6.8), Inches(1.8), Inches(0.5),
              t, font_size=12, color=C_WHITE, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# Slide 5: 10 节点流水线
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_bg(slide)
_add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
          "03  10 节点 Agent 流水线", font_size=32, color=C_ACCENT, bold=True)
_add_shape(slide, Inches(0.8), Inches(1.2), Inches(12), Inches(0.03), C_ACCENT)

agents = [
    ("① 协调", "coordinator", C_ACCENT),
    ("② 诊断", "diagnostician", C_ACCENT2),
    ("③ 规划", "planner", C_ACCENT3),
    ("④ 检索", "retriever", C_ACCENT),
    ("⑤ 生成", "generator", C_ACCENT2),
    ("⑥ 评估", "assessor", C_ACCENT3),
    ("⑦ 审阅", "critic", C_ACCENT),
    ("⑧ 路径", "path_planner", C_ACCENT2),
]
for i, (label, name, color) in enumerate(agents):
    x = Inches(0.5 + i * 1.55)
    y = Inches(1.8)
    _add_shape(slide, x, y, Inches(1.35), Inches(2.5), color)
    _add_shape(slide, x, y, Inches(1.35), Inches(0.06), color)
    _add_text(slide, x + Inches(0.1), y + Inches(0.3), Inches(1.15), Inches(0.5),
              label, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text(slide, x + Inches(0.1), y + Inches(0.9), Inches(1.15), Inches(0.4),
              name, font_size=11, color=C_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text(slide, x + Inches(0.1), y + Inches(1.4), Inches(1.15), Inches(0.9),
              ["解析请求", "分析薄弱", "制定计划", "检索知识", "并行生成", "质量评分", "事实核查", "输出路径"][i],
              font_size=11, color=C_LIGHT, alignment=PP_ALIGN.CENTER)

# 条件回跳箭头
_add_text(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.8),
          "④ 检索不到 → 回③重规划　　　⑦ 审阅不过 → 回④重检索+生成（最多 2 轮）",
          font_size=14, color=C_GRAY, alignment=PP_ALIGN.CENTER)

# generator_cluster 子 Agent
_add_text(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.4),
          "⑤ generator_cluster：7 个子 Agent 并行执行", font_size=15, color=C_ACCENT, bold=True)
subs = ["Teacher", "QuizMaster", "MindMap", "Extension", "Code", "PPT", "Video"]
for i, s in enumerate(subs):
    x = Inches(0.5 + i * 1.8)
    _add_shape(slide, x, Inches(6.0), Inches(1.6), Inches(0.5), C_ACCENT3)
    _add_text(slide, x, Inches(6.0), Inches(1.6), Inches(0.5),
              s, font_size=12, color=C_WHITE, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# Slide 6: GOMARL 共识
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_bg(slide)
_add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
          "04  GOMARL 共识机制", font_size=32, color=C_ACCENT, bold=True)
_add_shape(slide, Inches(0.8), Inches(1.2), Inches(12), Inches(0.03), C_ACCENT)

# 5 步流程
steps = [
    ("Step 1", "质量评分", "准确性/完整性/适配性\n1-10 分"),
    ("Step 2", "教学规则校验", "前置知识/难度\n适配性检查"),
    ("Step 3", "一致性校验", "语义矛盾检测\n事实冲突核查"),
    ("Step 4", "NeuralMixer", "PyTorch 神经网络\n动态权重融合"),
    ("Step 5", "决策", "通过/回炉/辩论\n三级判定"),
]
for i, (label, title, desc) in enumerate(steps):
    x = Inches(0.5 + i * 2.5)
    y = Inches(1.8)
    _add_shape(slide, x, y, Inches(2.2), Inches(3.5), C_ACCENT if i % 2 == 0 else C_ACCENT2)
    _add_shape(slide, x, y, Inches(2.2), Inches(0.06), C_ACCENT if i % 2 == 0 else C_ACCENT2)
    _add_text(slide, x + Inches(0.1), y + Inches(0.3), Inches(2.0), Inches(0.4),
              label, font_size=14, color=C_ACCENT if i % 2 == 0 else C_ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text(slide, x + Inches(0.1), y + Inches(0.8), Inches(2.0), Inches(0.5),
              title, font_size=18, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text(slide, x + Inches(0.1), y + Inches(1.5), Inches(2.0), Inches(1.5),
              desc, font_size=13, color=C_LIGHT, alignment=PP_ALIGN.CENTER)

_add_text(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
          "Agent 辩论：当 2 个 Agent 冲突时，启动 3 轮辩论协议 → 交叉质询 → 共识精炼",
          font_size=15, color=C_ACCENT3, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# Slide 7: FrugalRAG
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_bg(slide)
_add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
          "05  FrugalRAG 检索策略", font_size=32, color=C_ACCENT, bold=True)
_add_shape(slide, Inches(0.8), Inches(1.2), Inches(12), Inches(0.03), C_ACCENT)

# 7 步检索流程
retrieval_steps = [
    "① E5 向量化\n768 维嵌入",
    "② 向量检索\nTop-K 初筛",
    "③ 余弦过滤\n阈值 ≥ 0.65",
    "④ BM25 匹配\n关键词加权",
    "⑤ 融合排序\n0.7 向量 + 0.3 BM25",
    "⑥ 个性化重排\n5 因子画像调整",
    "⑦ Cross-encoder\n精排 Top-K",
]
for i, step in enumerate(retrieval_steps):
    x = Inches(0.3 + i * 1.85)
    y = Inches(1.8)
    _add_shape(slide, x, y, Inches(1.65), Inches(2.5), C_ACCENT2 if i % 2 == 0 else C_ACCENT3)
    _add_text(slide, x + Inches(0.1), y + Inches(0.2), Inches(1.45), Inches(2.1),
              step, font_size=12, color=C_WHITE, alignment=PP_ALIGN.CENTER)

# 个性化重排 5 因子
_add_text(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.4),
          "个性化重排 5 因子", font_size=18, color=C_ACCENT, bold=True)
factors = ["薄弱点 +0.15", "已掌握 -0.10", "考查权重 +0.10×w", "难度匹配 ±0.05", "高目标 +0.02"]
for i, f in enumerate(factors):
    x = Inches(0.5 + i * 2.5)
    _add_shape(slide, x, Inches(5.3), Inches(2.2), Inches(0.7), C_ACCENT)
    _add_text(slide, x, Inches(5.3), Inches(2.2), Inches(0.7),
              f, font_size=14, color=C_ACCENT, alignment=PP_ALIGN.CENTER)

_add_text(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
          "缓存策略：有 profile → 不缓存  |  无 profile → Redis 缓存 30 分钟",
          font_size=14, color=C_GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# Slide 8: 7 种资源
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_bg(slide)
_add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
          "06  7 种资源并行生成", font_size=32, color=C_ACCENT, bold=True)
_add_shape(slide, Inches(0.8), Inches(1.2), Inches(12), Inches(0.03), C_ACCENT)

resources = [
    ("📖 讲解文档", "Teacher Agent\n知识点详解", C_ACCENT),
    ("📝 练习题", "QuizMaster Agent\n含答案与解析", C_ACCENT2),
    ("🧠 思维导图", "MindMap Agent\n4步流水线生成", C_ACCENT3),
    ("📚 拓展阅读", "Extension Agent\n课外延伸材料", C_ACCENT),
    ("💻 代码实操", "CodePractice Agent\n可运行Python案例", C_ACCENT2),
    ("📊 PPT大纲", "PPTOutline Agent\n结构化幻灯片", C_ACCENT3),
    ("🎬 视频脚本", "VideoScript Agent\n讲解视频文案", C_ACCENT),
]
for i, (title, desc, color) in enumerate(resources):
    row = i // 4
    col = i % 4
    x = Inches(0.5 + col * 3.15)
    y = Inches(1.8 + row * 2.5)
    _add_shape(slide, x, y, Inches(2.9), Inches(2.1), color)
    _add_shape(slide, x, y, Inches(2.9), Inches(0.06), color)
    _add_text(slide, x + Inches(0.1), y + Inches(0.3), Inches(2.7), Inches(0.5),
              title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text(slide, x + Inches(0.1), y + Inches(1.0), Inches(2.7), Inches(0.8),
              desc, font_size=13, color=C_LIGHT, alignment=PP_ALIGN.CENTER)

_add_text(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
          "asyncio.gather() 并行执行 7 个 LLM 调用，return_exceptions=True 单个失败不影响整体",
          font_size=14, color=C_ACCENT3, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# Slide 9: NeuralMixer
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_bg(slide)
_add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
          "07  NeuralMixer 神经网络", font_size=32, color=C_ACCENT, bold=True)
_add_shape(slide, Inches(0.8), Inches(1.2), Inches(12), Inches(0.03), C_ACCENT)

# 三层结构
mixer_layers = [
    ("输入层", "E5 编码器 → 768 维向量\n7 个 Agent 输出文本 → 语义向量", C_ACCENT2, Inches(1.6)),
    ("权重层", "EWMA 历史表现 × 学生画像 × 教学规则\n动态权重计算，归一化后参与融合", C_ACCENT, Inches(3.0)),
    ("决策层", "GroupMixer 神经网络\n多头注意力 + 组内相似度 + 组间多样性 + Lasso 正则", C_ACCENT3, Inches(4.4)),
]
for name, desc, color, y in mixer_layers:
    _add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.2), color)
    _add_shape(slide, Inches(0.8), y, Inches(0.08), Inches(1.2), color)
    _add_text(slide, Inches(1.2), y + Inches(0.1), Inches(3), Inches(0.4),
              name, font_size=20, color=color, bold=True)
    _add_text(slide, Inches(1.2), y + Inches(0.5), Inches(10), Inches(0.6),
              desc, font_size=14, color=C_LIGHT)

_add_text(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.8),
          "降级策略：PyTorch 不可用时 → 加权平均（规则模式）\n"
          "训练策略：500 条 seed data 样本，同类知识点高相似度 → 高共识分数",
          font_size=14, color=C_GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# Slide 10: 总结
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_bg(slide)
_add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
          "08  总结与展望", font_size=32, color=C_ACCENT, bold=True)
_add_shape(slide, Inches(0.8), Inches(1.2), Inches(12), Inches(0.03), C_ACCENT)

# 优势
_add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
          "系统优势", font_size=22, color=C_ACCENT, bold=True)
_add_bullet_text(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(3), [
    "多智能体协同：10 节点流水线 + 7 子 Agent 并行",
    "质量保障：GOMARL 评分 + Critic 审阅 + 辩论机制",
    "个性化：8 维度画像 + MindMap 掌握度标注",
    "容错：LLM 降级 + 重试 2 轮 + 强制通过",
])

# 未来
_add_text(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.5),
          "未来方向", font_size=22, color=C_ACCENT2, bold=True)
_add_bullet_text(slide, Inches(6.8), Inches(2.2), Inches(5.5), Inches(3), [
    "FrugalRAG SFT 微调 + GRPO 强化",
    "GOMARL 证据冲突消解完整版",
    "Milvus 生产部署 + 四科全覆盖",
    "多模态：讯飞 TTI/TTS/PPT/数字人",
])

# 数据
_add_text(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.8),
          "后端：281 测试全绿  |  前端：30 组件零错误  |  知识库：571 chunks + 200 题  |  安全：21/21 项修复",
          font_size=14, color=C_ACCENT3, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# Slide 11: 问答
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_bg(slide)
_add_text(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(1),
          "Q & A", font_size=60, color=C_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
_add_text(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.5),
          "感谢聆听！欢迎提问", font_size=24, color=C_WHITE, alignment=PP_ALIGN.CENTER)
_add_shape(slide, Inches(5), Inches(4.8), Inches(3.333), Inches(0.06), C_ACCENT)
_add_text(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5),
          "基于 GOMARL 与 FrugalRAG 的 408 考研个性化学习多智能体系统",
          font_size=16, color=C_GRAY, alignment=PP_ALIGN.CENTER)

# ── 保存 ──
output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "MARS-408_软件杯演示.pptx")
prs.save(output_path)
print(f"✅ PPT 已保存: {output_path}")
print(f"   共 {len(prs.slides)} 页幻灯片")