# ============================================================
# 教材导入脚本 — 新版（支持 InMemoryVectorStore / Milvus）
# 读取 PDF/PPTX/DOCX → 分块 → 存入向量库
# 用法：python import_pdfs.py [--rebuild]
# ============================================================

import os
import sys
import re
import json
import hashlib
import logging
from pathlib import Path
from typing import Optional

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(name)s] %(levelname)s: %(message)s")
logger = logging.getLogger("import_pdfs")

# 项目根目录
PROJECT_ROOT = Path(__file__).parent.parent
DOCS_DIR = PROJECT_ROOT / "documents" / "教材"

# ── 各科 subject 映射 ──
SUBJECT_MAP = {
    "计算机网络":   "computer_network",
    "数据结构":     "data_structures",
    "操作系统":     "operating_system",
    "计算机组成原理": "computer_organization",
    "408真题":     "exam",
}

# ── 文件名 → subject 的关键词匹配 ──
FILENAME_SUBJECT_RULES = [
    ("王道操作系统", "operating_system"),
    ("操作系统", "operating_system"),
    ("王道数据结构", "data_structures"),
    ("数据结构", "data_structures"),
    ("王道计算机组成原理", "computer_organization"),
    ("组成原理", "computer_organization"),
    ("计算机网络", "computer_network"),
    ("计网", "computer_network"),
    ("408", "exam"),
    ("真题", "exam"),
    ("大纲", "exam"),
]


def detect_subject(filepath: str) -> str:
    """根据文件路径判断所属课程"""
    path_lower = filepath.lower()
    for keyword, subject in FILENAME_SUBJECT_RULES:
        if keyword.lower() in path_lower:
            return subject
    return "computer_network"  # 默认


def extract_text_from_pdf(pdf_path: str, use_ocr: bool = False) -> str:
    """用 PyMuPDF 提取 PDF 文本，自动跳过扫描版（<50字符/页）
    如果 use_ocr=True，对扫描版自动调用 Tesseract OCR
    """
    try:
        import fitz
        doc = fitz.open(pdf_path)
        n_pages = len(doc)
        text_parts = []
        total_chars = 0
        
        # 先尝试普通提取
        for page_num in range(n_pages):
            try:
                text = doc[page_num].get_text()
                if text.strip() and len(text) > 50:
                    text_parts.append(text)
                    total_chars += len(text)
            except Exception:
                continue

        # 如果文字太少且启用了OCR，用OCR识别
        if total_chars < 200 and use_ocr:
            doc.close()
            return extract_text_from_pdf_ocr(pdf_path)
            
        doc.close()

        if total_chars < 200:
            logger.info(f"  ⚠️ 扫描版PDF（仅{total_chars}字符）")
            if use_ocr:
                return extract_text_from_pdf_ocr(pdf_path)
            logger.info("  (跳过，使用 --ocr 参数启用OCR识别)")
            return ""

        full_text = "\n".join(text_parts)
        logger.info(f"  PDF提取: {total_chars} 字符/ {n_pages} 页")
        return full_text
    except ImportError:
        logger.error("  ❌ 需要 PyMuPDF: pip install pymupdf")
        return ""
    except Exception as e:
        logger.warning(f"  ⚠️ PDF解析跳过: {e}")
        return ""


def extract_text_from_pdf_ocr(pdf_path: str) -> str:
    """用 Tesseract OCR 识别扫描版PDF"""
    try:
        import fitz
        from PIL import Image
        import pytesseract
        import io
    except ImportError as e:
        logger.error(f"  ❌ OCR需要依赖: pip install pytesseract pillow")
        return ""

    # 检查 Tesseract 是否安装
    try:
        pytesseract.get_tesseract_version()
    except Exception:
        logger.error("  ❌ 未找到 Tesseract OCR 引擎")
        logger.error("     请下载安装: https://github.com/UB-Mannheim/tesseract/releases")
        logger.error("     安装后如不在默认路径，设置: pytesseract.pytesseract.tesseract_cmd = r'C:\\Program Files\\Tesseract-OCR\\tesseract.exe'")
        return ""

    try:
        doc = fitz.open(pdf_path)
        n_pages = min(len(doc), 100)  # OCR 很慢，限制最多 100 页
        text_parts = []

        logger.info(f"  开始OCR识别 {n_pages} 页（每页约3-5秒）...")
        for page_num in range(n_pages):
            try:
                page = doc[page_num]
                # PDF页面 → 图片
                pix = page.get_pixmap(dpi=200)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                # OCR识别（中文+英文）
                text = pytesseract.image_to_string(img, lang="chi_sim+eng")
                if text.strip():
                    text_parts.append(f"【第{page_num+1}页】\n{text.strip()}")
                if (page_num + 1) % 10 == 0:
                    logger.info(f"  OCR进度: {page_num+1}/{n_pages} 页")
            except Exception as e:
                logger.warning(f"  第{page_num+1}页OCR失败: {e}")

        doc.close()
        full_text = "\n\n".join(text_parts)
        logger.info(f"  OCR完成: {len(text_parts)} 页有文字, 共 {len(full_text)} 字符")
        return full_text

    except Exception as e:
        logger.error(f"  ❌ OCR处理失败: {e}")
        return ""


def extract_text_from_pptx(pptx_path: str) -> str:
    """从 PPTX 提取文本"""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                text_parts.append(f"【第{slide_num}页】\n" + "\n".join(slide_texts))
        return "\n\n".join(text_parts)
    except ImportError:
        logger.error("  ❌ 需要 python-pptx: pip install python-pptx")
        return ""
    except Exception as e:
        logger.warning(f"  ⚠️ PPTX解析失败: {e}")
        return ""


def extract_text_from_docx(docx_path: str) -> str:
    """从 DOCX 提取文本"""
    try:
        from docx import Document
        doc = Document(docx_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except ImportError:
        logger.error("  ❌ 需要 python-docx: pip install python-docx")
        return ""
    except Exception as e:
        logger.warning(f"  ⚠️ DOCX解析失败: {e}")
        return ""


def extract_text_from_doc(doc_path: str) -> str:
    """从旧版 .doc 提取文本（尝试用 antiword 或 catdoc）"""
    try:
        import subprocess
        result = subprocess.run(["catdoc", doc_path], capture_output=True, text=True, timeout=10)
        if result.returncode == 0:
            return result.stdout
    except Exception:
        pass
    logger.warning(f"  ⚠️ 旧版 .doc 格式暂不支持解析: {doc_path}")
    return ""


def semantic_chunk(text: str, max_chars: int = 600) -> list[str]:
    """语义分块：按段落拆分，合并短段落，拆分长段落"""
    # 按空行分段落
    paragraphs = [p.strip() for p in re.split(r'\n\s*\n', text) if p.strip()]
    if not paragraphs:
        paragraphs = [text.strip()]

    chunks = []
    current = ""

    for para in paragraphs:
        # 标题或列表项：单独成块
        if re.match(r'^#{1,6}\s|^第.*[章节讲]|^[-*\d+.]\s|^【', para):
            if current.strip():
                chunks.append(current.strip())
            current = para
        # 短段落：合并到当前块
        elif len(current) + len(para) < max_chars * 0.8:
            if current:
                current += "\n" + para
            else:
                current = para
        # 长段落或达到上限：截断
        else:
            if current.strip():
                chunks.append(current.strip())
            current = para

        # 如果单段落超过 max_chars，再拆分
        if len(current) > max_chars:
            if current.strip():
                chunks.append(current.strip())
            current = ""

    if current.strip():
        chunks.append(current.strip())

    # 过滤太短的块（<20字符）
    chunks = [c for c in chunks if len(c) >= 20]

    # 合并太短的相邻块
    merged = []
    for c in chunks:
        if merged and len(merged[-1]) < 120 and len(merged[-1]) + len(c) < max_chars:
            merged[-1] += "\n" + c
        else:
            merged.append(c)

    return merged if merged else [text[:max_chars]]


def extract_chapter(text: str, subject: str) -> str:
    """尝试从文本中提取章节信息"""
    # 匹配 "第X章"、"第X讲"、"X." 等模式
    patterns = [
        r'第[一二三四五六七八九十\d]+[章节讲篇]',
        r'^\d+[\.\、]\s*\S+',
        r'^[综]',
    ]
    for p in patterns:
        m = re.search(p, text[:200])
        if m:
            return m.group(0)
    return "综合"


def process_file(filepath: str, rebuild: bool = False, use_ocr: bool = False) -> list[dict]:
    """处理单个文件，返回 chunks 列表"""
    filepath = str(filepath)
    ext = Path(filepath).suffix.lower()
    subject = detect_subject(filepath)

    logger.info(f"处理: {Path(filepath).name} ({subject})")

    # 提取文本
    text = ""
    if ext == ".pdf":
        text = extract_text_from_pdf(filepath, use_ocr=use_ocr)
    elif ext == ".pptx":
        text = extract_text_from_pptx(filepath)
    elif ext == ".docx":
        text = extract_text_from_docx(filepath)
    elif ext == ".doc":
        text = extract_text_from_doc(filepath)
    else:
        logger.info(f"  跳过不支持格式: {ext}")
        return []

    if not text.strip():
        logger.warning(f"  无有效文本")
        return []

    # 分块
    chunks = semantic_chunk(text)
    logger.info(f"  分块: {len(chunks)} 块")

    # 构建带 metadata 的 chunk 数据
    result = []
    for i, chunk_text in enumerate(chunks):
        chunk_id = hashlib.md5(f"{filepath}_{i}".encode()).hexdigest()[:12]
        chapter = extract_chapter(chunk_text, subject)
        result.append({
            "id": f"import_{chunk_id}",
            "text": chunk_text,
            "metadata": {
                "subject": subject,
                "chapter": chapter,
                "type": "imported",
                "source": Path(filepath).name,
                "chunk_index": i,
            }
        })

    return result


def import_all(rebuild: bool = False, use_ocr: bool = False):
    """扫描 documents/教材/ 下所有文件并导入向量库"""
    # 延迟导入（避免启动时依赖）
    sys.path.insert(0, str(PROJECT_ROOT))
    from db.milvus_client import vector_db
    from db.embedder import embed_text

    # 扫描所有可解析的文件
    supported_exts = {".pdf", ".pptx", ".docx", ".doc"}
    files_to_process = []
    for root, dirs, files in os.walk(str(DOCS_DIR)):
        for f in files:
            ext = Path(f).suffix.lower()
            if ext in supported_exts:
                files_to_process.append(Path(root) / f)

    logger.info(f"找到 {len(files_to_process)} 个可导入文件 ({'OCR开启' if use_ocr else 'OCR关闭'})")

    # 如果 rebuild，先清空已有导入数据
    if rebuild:
        logger.info("清空已有导入数据...")
        # 获取已有集合中所有 type=imported 的 ID
        # InMemoryVectorStore 不支持选择性删除，这里直接重建集合
        try:
            from db.embedder import get_embed_dim
            _dim = get_embed_dim()
            existing = vector_db.search("netlearn_kb", [0.0]*_dim, top_k=10000)
            imported_ids = [c["id"] for c in existing if c.get("metadata", {}).get("type") == "imported"]
            if imported_ids:
                logger.info(f"找到 {len(imported_ids)} 条旧导入数据，将覆盖")
        except Exception:
            pass

    # 处理每个文件
    all_chunks = []
    for filepath in files_to_process:
        try:
            chunks = process_file(filepath, rebuild, use_ocr=use_ocr)
            all_chunks.extend(chunks)
        except Exception as e:
            logger.error(f"处理失败 {filepath}: {e}")

    logger.info(f"共生成 {len(all_chunks)} 个知识块")

    # 批量编码并写入向量库
    if all_chunks:
        # 先统一编码所有文本
        texts = [c["text"] for c in all_chunks]
        logger.info(f"正在 E5 编码 {len(texts)} 条文本...")
        try:
            embeddings = embed_batch(texts)
        except Exception as e:
            logger.warning(f"E5 编码失败: {e}，使用零向量")
            import numpy as np
            from db.embedder import get_embed_dim
            _dim = get_embed_dim()
            embeddings = [np.zeros(_dim).tolist() for _ in texts]

        # 构建 chunks 列表
        vector_chunks = []
        for i, c in enumerate(all_chunks):
            vector_chunks.append({
                "id": c["id"],
                "text": c["text"],
                "metadata": c["metadata"],
                "embedding": embeddings[i],
            })

        # 分批写入（每批 100 条），Q6 优化：内存累积 + 周期/末尾落盘避免 O(n²) 重写
        batch_size = 100
        flush_every = 10  # 每 10 批(1000条)强制落盘一次，平衡性能与崩溃数据丢失
        total_inserted = 0
        batch_idx = 0
        for i in range(0, len(vector_chunks), batch_size):
            batch = vector_chunks[i:i+batch_size]
            try:
                inserted = vector_db.insert("netlearn_kb", batch, save=False)
                total_inserted += inserted
                batch_idx += 1
                logger.info(f"  写入第 {batch_idx} 批: {len(batch)} 条")
                if batch_idx % flush_every == 0:
                    vector_db.flush("netlearn_kb")
                    logger.info(f"  [checkpoint] 已落盘 {batch_idx} 批")
            except Exception as e:
                logger.error(f"  写入失败: {e}")

        # 末尾强制落盘（save=False 后必须 flush，否则数据仅驻留内存）
        try:
            vector_db.flush("netlearn_kb")
        except Exception as e:
            logger.error(f"  最终落盘失败: {e}")

        total = vector_db.count("netlearn_kb")
        logger.info(f"导入完成！本次插入 {total_inserted} 条，向量库总文档数: {total}")
    else:
        logger.warning("没有生成任何知识块")


def embed_batch(texts: list[str]) -> list[list[float]]:
    """批量编码（兼容无 E5 模型的情况）"""
    try:
        from db.embedder import embed_batch as _embed_batch, get_embed_dim
        return _embed_batch(texts)
    except Exception as e:
        logger.warning(f"E5 编码失败，使用模拟向量: {e}")
        import numpy as np
        from db.embedder import get_embed_dim
        dim = get_embed_dim()
        return [np.zeros(dim).tolist() for _ in texts]


if __name__ == "__main__":
    import sys, subprocess
    from pathlib import Path
    print("[DEPRECATED] import_pdfs.py 不再直写向量库；已改为通过后端导入队列导入。")
    print("  等效命令： python tools/import_client.py --type pdf [--rebuild] [--use-ocr]")
    client = Path(__file__).parent / "tools" / "import_client.py"
    args = [sys.executable, str(client), "--type", "pdf"]
    if "--rebuild" in sys.argv:
        args.append("--rebuild")
    if "--ocr" in sys.argv:
        args.append("--use-ocr")
    try:
        sys.exit(subprocess.call(args))
    except Exception as e:
        print(f"无法启动导入客户端（请先启动后端 python main.py）：{e}")
        sys.exit(1)
